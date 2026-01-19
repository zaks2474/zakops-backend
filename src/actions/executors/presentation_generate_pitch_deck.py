from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, Optional

from actions.engine.models import ActionError, ActionPayload, ArtifactMetadata, now_utc_iso
from actions.executors._artifacts import resolve_action_artifact_dir
from actions.executors.base import ActionExecutionError, ActionExecutor, ExecutionContext, ExecutionResult


class GeneratePitchDeckExecutor(ActionExecutor):
    action_type = "PRESENTATION.GENERATE_PITCH_DECK"

    def validate(self, payload: ActionPayload) -> tuple[bool, Optional[str]]:
        return True, None

    def execute(self, payload: ActionPayload, ctx: ExecutionContext) -> ExecutionResult:
        inputs = payload.inputs or {}
        deck_title = str(inputs.get("deck_title") or (ctx.deal or {}).get("canonical_name") or "Pitch Deck").strip()
        subtitle = str(inputs.get("subtitle") or "").strip()
        include_financials = bool(inputs.get("include_financials", True))

        out_dir = resolve_action_artifact_dir(ctx)

        # Required PPTX: use python-pptx if available (clear failure if missing).
        try:
            from pptx import Presentation  # type: ignore
        except Exception as e:
            raise ActionExecutionError(
                ActionError(
                    code="python_pptx_missing",
                    message="python-pptx is required to generate PPTX pitch decks (pip install python-pptx).",
                    category="dependency",
                    retryable=False,
                    details={"import_error": type(e).__name__},
                )
            )

        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = deck_title
        if subtitle and slide.placeholders and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle

        # Simple overview slide
        bullet_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_layout)
        slide2.shapes.title.text = "Overview"
        body = slide2.shapes.placeholders[1].text_frame
        body.text = "Deal summary"
        body.add_paragraph().text = "Key highlights"
        if include_financials:
            body.add_paragraph().text = "Financial snapshot"

        pptx_path = out_dir / "pitch_deck.pptx"
        prs.save(str(pptx_path))

        artifacts = [
            ArtifactMetadata(
                filename=pptx_path.name,
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                path=str(pptx_path),
                created_at=now_utc_iso(),
            )
        ]

        outputs: Dict[str, Any] = {"deck_title": deck_title, "include_financials": include_financials}
        return ExecutionResult(outputs=outputs, artifacts=artifacts)

